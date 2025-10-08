from smolagents import Tool
import json
import os
from typing import Dict, List
from datetime import datetime

class ReportGeneratorTool(Tool):
    name = "report_generator"
    description = "Generate PowerPoint presentations from final analysis results and stored visualizations. Use the analysis JSON output and image files from plots/ directory to create professional slides with embedded charts and insights."
    inputs = {
        "analysis_results": {
            "type": "object",
            "description": "Dictionary containing the analysis results and insights"
        },
        "plots": {
            "type": "object",
            "description": "List of plot file paths to include in the presentation (optional - will auto-detect from plots/ directory)",
            "nullable": True
        }
    }
    output_type = "string"

    def forward(self, analysis_results: Dict, plots: List[str] = None) -> str:
        """Generate PowerPoint presentation from analysis results"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            # If no plots provided, automatically find PNG files in plots directory
            if plots is None or len(plots) == 0:
                plots = []
                if os.path.exists('plots'):
                    plot_files = [f for f in os.listdir('plots') if f.endswith('.png')]
                    plots = [f'plots/{f}' for f in plot_files]
                    print(f"📊 Found {len(plots)} PNG files in plots/ directory")

            # Create presentation
            prs = Presentation()

            # Helper function to add title slide
            def add_title_slide(title, subtitle=""):
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title_box = slide.shapes.title
                title_box.text = title
                title_box.text_frame.paragraphs[0].font.size = Pt(32)  # Reduced from 44
                title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 160)
                title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                if subtitle:
                    subtitle_box = slide.placeholders[1]
                    subtitle_box.text = subtitle
                    subtitle_box.text_frame.paragraphs[0].font.size = Pt(18)  # Reduced from 24
                    subtitle_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)

            # Helper function to add content slide
            def add_content_slide(title, content_list):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title_box = slide.shapes.title
                title_box.text = title
                title_box.text_frame.paragraphs[0].font.size = Pt(28)  # Reduced from 36
                title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 160)

                content_box = slide.placeholders[1]
                content_box.text = ""
                for item in content_list:
                    p = content_box.text_frame.add_paragraph()
                    p.text = str(item)
                    p.font.size = Pt(16)  # Reduced from 20
                    p.level = 0

            # Helper function to add slide with image
            def add_image_slide(title, image_path, explanation=""):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title_box = slide.shapes.title
                title_box.text = title
                title_box.text_frame.paragraphs[0].font.size = Pt(24)  # Reduced from 32
                title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 160)

                if os.path.exists(image_path):
                    # Add image
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(8)
                    height = Inches(4.5)
                    slide.shapes.add_picture(image_path, left, top, width, height)

                    # Add explanation below image
                    if explanation:
                        left = Inches(1)
                        top = Inches(6.5)
                        width = Inches(8)
                        height = Inches(1)
                        textbox = slide.shapes.add_textbox(left, top, width, height)
                        textbox.text = explanation
                        textbox.text_frame.paragraphs[0].font.size = Pt(14)  # Reduced from 16
                        textbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)
                else:
                    content_box = slide.placeholders[1]
                    content_box.text = f"Image not found: {image_path}"

            # 1. Title Slide
            title = analysis_results.get('title', 'Data Analysis Report')
            current_date = datetime.now().strftime("%B %d, %Y")
            add_title_slide(title, f"Generated on {current_date}")

            # 2. Dataset Overview Slide
            dataset_info = []
            if 'dataset_overview' in analysis_results:
                overview = analysis_results['dataset_overview']
                if 'shape' in overview:
                    dataset_info.append(f"Dataset Shape: {overview['shape']}")
                if 'columns' in overview:
                    dataset_info.append(f"Number of Columns: {len(overview['columns'])}")
                if 'sample_rows' in overview:
                    dataset_info.append(f"Sample data available: {len(overview['sample_rows'])} rows shown")
            
            if 'eda_summary' in analysis_results:
                eda = analysis_results['eda_summary']
                for key, value in eda.items():
                    if isinstance(value, dict) and value:
                        # Extract meaningful info from EDA summary
                        if 'top_productlines' in key.lower():
                            top_item = list(value.keys())[0] if value else "N/A"
                            dataset_info.append(f"Top Product Line: {top_item}")
                        elif 'dealsize' in key.lower():
                            total_deals = sum(value.values()) if isinstance(value, dict) else "N/A"
                            dataset_info.append(f"Total Deals: {total_deals}")
            
            if not dataset_info:
                dataset_info = ["Dataset successfully loaded and analyzed", "Comprehensive EDA performed"]
            
            add_content_slide("Dataset Overview", dataset_info)

            # 3. Data Quality Slide
            data_quality_points = []
            if 'analysis_sections' in analysis_results:
                for section in analysis_results['analysis_sections']:
                    if section.get('type') == 'subsection_title' and 'Data Quality' in section.get('content', ''):
                        # Find the next paragraph
                        section_index = analysis_results['analysis_sections'].index(section)
                        if section_index + 1 < len(analysis_results['analysis_sections']):
                            next_section = analysis_results['analysis_sections'][section_index + 1]
                            if next_section.get('type') == 'paragraph':
                                data_quality_points.append(next_section['content'])
                        break

            if not data_quality_points:
                data_quality_points = ["Data quality assessment completed", "Basic data validation performed"]

            add_content_slide("Data Quality", data_quality_points)

            # 4. Add ALL Visualization Slides from detected plots
            if plots:
                print(f"🎯 Creating slides for {len(plots)} visualizations...")
                
                # Define meaningful titles and descriptions for common plot types
                plot_info = {
                    'sales_hist': ('Sales Distribution', 'Histogram showing the distribution of sales values across the dataset'),
                    'sales_by_productline': ('Sales by Product Line', 'Bar chart comparing total sales across different product lines'),
                    'sales_by_country': ('Sales by Country', 'Top performing countries by total sales volume'),
                    'sales_trend': ('Sales Trend Analysis', 'Time series analysis showing sales patterns over time'),
                    'corr_heatmap': ('Correlation Heatmap', 'Correlation matrix showing relationships between numeric variables'),
                    'price_vs_quantity': ('Price vs Quantity Analysis', 'Scatter plot exploring the relationship between price and quantity'),
                    'boxplot_sales': ('Sales Distribution by Category', 'Box plot showing sales distribution across different categories'),
                    'top_customers': ('Top Customers Analysis', 'Analysis of highest value customers by sales volume'),
                    'model_performance': ('Model Performance Comparison', 'Comparison of different machine learning model performances'),
                    'feature_importances': ('Feature Importance Analysis', 'Analysis showing which features are most important for predictions'),
                    'rf_feature': ('Random Forest Feature Importance', 'Random Forest model feature importance ranking'),
                    'gbr_feature': ('Gradient Boosting Feature Importance', 'Gradient Boosting model feature importance ranking')
                }
                
                for i, plot_path in enumerate(plots):
                    if os.path.exists(plot_path):
                        # Extract plot type from filename
                        filename = os.path.basename(plot_path).lower().replace('.png', '')
                        
                        # Find matching plot info
                        title = f"Visualization {i+1}"
                        explanation = f"Analysis visualization from {filename}"
                        
                        for key, (plot_title, plot_desc) in plot_info.items():
                            if key in filename:
                                title = plot_title
                                explanation = plot_desc
                                break
                        
                        # Use plot_descriptions from analysis_results if available
                        if 'plot_descriptions' in analysis_results:
                            plot_key = plot_path.replace('\\', '/')
                            if plot_key in analysis_results['plot_descriptions']:
                                explanation = analysis_results['plot_descriptions'][plot_key]
                        
                        print(f"📊 Adding slide {i+1}: {title}")
                        add_image_slide(title, plot_path, explanation)
                    else:
                        print(f"⚠️ Plot file not found: {plot_path}")

            # 5. Legacy support for visualization_sections (if any exist)
            visualization_sections = []
            if 'analysis_sections' in analysis_results:
                for section in analysis_results['analysis_sections']:
                    if section.get('type') == 'visualization':
                        visualization_sections.append(section)

            for i, viz_section in enumerate(visualization_sections):
                image_path = viz_section['content']['image_path']
                if image_path not in plots:  # Only add if not already added above
                    explanation = viz_section['content']['explanation']
                    title = f"Additional Visualization {i+1}"
                    add_image_slide(title, image_path, explanation)

            # 5. Key Insights Slide
            insights = []
            
            # Extract insights from various sources
            if 'key_findings' in analysis_results:
                findings = analysis_results['key_findings']
                for key, value in findings.items():
                    if isinstance(value, str):
                        insights.append(f"{key.replace('_', ' ').title()}: {value}")
                    elif isinstance(value, dict):
                        insights.append(f"{key.replace('_', ' ').title()}: {str(value)}")
            
            if 'eda_summary' in analysis_results:
                eda = analysis_results['eda_summary']
                for key, value in eda.items():
                    if isinstance(value, dict) and value:
                        top_item = list(value.keys())[0] if value else None
                        if top_item:
                            insights.append(f"Top {key.replace('_', ' ')}: {top_item}")
            
            # Extract from analysis_sections if available
            if 'analysis_sections' in analysis_results:
                for section in analysis_results['analysis_sections']:
                    if section.get('type') == 'list':
                        insights.extend(section.get('content', []))
                    elif section.get('type') == 'paragraph':
                        content = section.get('content', '')
                        if len(content) < 200 and 'insight' in content.lower():
                            insights.append(content)
            
            # Add model-related insights
            if 'model' in analysis_results:
                model_info = analysis_results['model']
                if 'task' in model_info:
                    insights.append(f"ML Task: {model_info['task']}")
                if 'top_features_sample' in model_info:
                    features = model_info['top_features_sample']
                    if features:
                        top_feature = features[0] if isinstance(features[0], dict) else str(features[0])
                        if isinstance(top_feature, dict):
                            feature_name = top_feature.get('feature', 'Unknown')
                            insights.append(f"Most important feature: {feature_name}")

            if not insights:
                insights = ["Comprehensive data analysis completed", "Multiple visualizations generated", "Statistical patterns identified"]

            add_content_slide("Key Insights", insights[:8])  # Limit to 8 insights

            # 6. Machine Learning Results Slide
            if 'model' in analysis_results:
                model_info = analysis_results['model']
                ml_content = []
                
                if 'task' in model_info:
                    ml_content.append(f"Task: {model_info['task']}")
                
                if 'artifacts' in model_info:
                    artifacts = model_info['artifacts']
                    artifact_count = sum(1 for v in artifacts.values() if v and os.path.exists(str(v)))
                    ml_content.append(f"Generated {artifact_count} model artifacts")
                
                if 'top_features_sample' in model_info:
                    features = model_info['top_features_sample']
                    if features:
                        ml_content.append(f"Analyzed {len(features)} feature importances")
                        if len(features) > 0:
                            top_3 = features[:3]
                            ml_content.append("Top 3 Most Important Features:")
                            for i, feature in enumerate(top_3):
                                if isinstance(feature, dict):
                                    name = feature.get('feature', f'Feature {i+1}')
                                    importance = feature.get('importance', 'N/A')
                                    ml_content.append(f"  {i+1}. {name}: {importance}")
                
                if 'notes' in model_info:
                    ml_content.append(f"Notes: {model_info['notes']}")
                
                if ml_content:
                    add_content_slide("Machine Learning Results", ml_content)

            # 7. Recommendations Slide
            recommendations_content = []
            if 'recommendations' in analysis_results:
                recommendations_content = analysis_results['recommendations']
            else:
                # Generate default recommendations based on available data
                recommendations_content = [
                    "Review the generated visualizations for actionable insights",
                    "Consider the machine learning model results for predictive capabilities",
                    "Focus on the top-performing categories identified in the analysis",
                    "Investigate any anomalies or outliers discovered in the data"
                ]
            
            add_content_slide("Recommendations", recommendations_content)

            # 8. Conclusion Slide
            conclusion_content = []
            if 'conclusion' in analysis_results:
                conclusion_content.append(analysis_results['conclusion'])
            else:
                conclusion_content.append("Comprehensive data analysis completed successfully")
                if plots:
                    conclusion_content.append(f"Generated {len(plots)} visualizations for detailed insights")
                if 'model' in analysis_results:
                    conclusion_content.append("Machine learning models built and evaluated")

            conclusion_content.append("\nThank you for using Spark Insights AI Data Scientist Assistant!")

            add_content_slide("Conclusion", conclusion_content)

            # Save the presentation
            output_file = "analysis_report.pptx"
            prs.save(output_file)

            file_size = os.path.getsize(output_file)
            return f"PowerPoint presentation created successfully!\nFile: {output_file}\nSize: {file_size} bytes\nSlides: {len(prs.slides)}"

        except ImportError:
            return "Error: python-pptx library not installed. Please install with: pip install python-pptx"
        except Exception as e:
            return f"Error generating PowerPoint presentation: {str(e)}"